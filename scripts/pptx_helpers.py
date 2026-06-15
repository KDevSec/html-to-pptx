#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Reusable python-pptx helpers for HTML->PPTX replication.

Coordinate model (see references/playbook.md): capture the page at 1280x720 / dsf=1,
then 1 css px = 9525 EMU and font pt = css px * 0.75. Every getBoundingClientRect()
value maps straight onto the 16:9 deck, so native text / shapes / image slices that
share one capture line up exactly.

Import this from a per-figure build script:
    import sys; sys.path.insert(0, '<skill>/scripts')
    from pptx_helpers import *
    prs, slide = new_deck()
    bg_fill(slide, '#05070c')
    rrect(slide, x, y, w, h, 12, fill='#ffffff', fa=8, line='#ffffff', la=20)
    textbox(slide, x, y, w, h, [[('hi', P(13), rgb('#eaf2ff'), True)]])
    picture(slide, 'icon.png', ix, iy, iw, ih)
    save(prs, 'out.pptx')
"""
import re, zipfile
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

EMU_W, EMU_H = 12192000, 6858000   # 16:9 13.333 x 7.5 in
PXE = 9525                         # EMU per css px at 1280-wide capture
_FONT = 'Microsoft YaHei'         # 微软雅黑 (default — ships with Windows); change via set_font_family()/target_cjk_font()

def set_font_family(name):
    global _FONT; _FONT = name

# CJK font that exists on the machine where the deck will be PRESENTED (not where
# it's generated) — a missing font triggers substitution and shifts layout. Pick
# by the audience's OS, or embed fonts for true portability (see playbook §compat).
CJK_FONT = {
    'windows': 'Microsoft YaHei',   # 微软雅黑 — ships with Windows
    'mac':     'PingFang SC',        # ships with macOS
    'linux':   'Noto Sans CJK SC',   # common on Linux
    'office365': 'Microsoft YaHei',
}
def target_cjk_font(target='windows'):
    """Set the font for the audience's platform. Default Windows (most common for
    Office). For cross-platform delivery, prefer embedding fonts in PowerPoint."""
    set_font_family(CJK_FONT.get(target, 'Microsoft YaHei'))
    return _FONT

def E(px):  return Emu(int(round(px * PXE)))
def P(px):  return Pt(round(px * 0.75, 2))      # css px -> pt

def rgb(s, default=(0xEA, 0xF2, 0xFF)):
    """Accept '#abc', '#aabbcc', 'rgb(r,g,b)', 'rgba(r,g,b,a)', or an RGBColor."""
    if isinstance(s, RGBColor): return s
    if not s: return RGBColor(*default)
    s = s.strip()
    if s.startswith('#'):
        h = s[1:]
        if len(h) == 3: h = ''.join(c * 2 for c in h)
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    m = re.findall(r'\d+\.?\d*', s)
    if len(m) >= 3: return RGBColor(int(float(m[0])), int(float(m[1])), int(float(m[2])))
    return RGBColor(*default)

def set_run(run, size, color, bold=False, font=None):
    font = font or _FONT
    f = run.font; f.size = size; f.bold = bold; f.color.rgb = (color if isinstance(color, RGBColor) else rgb(color)); f.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):     # CJK needs all three typeface slots
        e = rPr.find(qn(tag))
        if e is None: e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set('typeface', font)

def new_deck():
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
    return prs, slide

# ---- alpha / dash low-level (python-pptx has no direct API) ----
def _alpha(sp, fa=None, la=None):
    spPr = sp._element.spPr
    if fa is not None:
        sf = spPr.find(qn('a:solidFill'))
        if sf is not None:
            c = sf.find(qn('a:srgbClr'))
            if c is not None: c.append(c.makeelement(qn('a:alpha'), {'val': str(int(fa * 1000))}))
    if la is not None:
        ln = spPr.find(qn('a:ln'))
        if ln is not None:
            sf = ln.find(qn('a:solidFill'))
            if sf is not None:
                c = sf.find(qn('a:srgbClr'))
                if c is not None: c.append(c.makeelement(qn('a:alpha'), {'val': str(int(la * 1000))}))

def _dash(sp, val='dash'):
    ln = sp._element.spPr.find(qn('a:ln'))
    if ln is not None: ln.append(ln.makeelement(qn('a:prstDash'), {'val': val}))

def rrect(slide, x, y, w, h, rad_px=10, fill=None, fa=None, line=None, la=None, lw=1.0, dash=None):
    """Rounded rectangle. fill/line accept color str or None (none). fa/la = alpha 0-100."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, E(x), E(y), E(max(w, 2)), E(max(h, 2)))
    try: sp.adjustments[0] = min(0.5, rad_px / max(1, min(w, h)))
    except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = rgb(line); sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    _alpha(sp, fa, la)
    if dash: _dash(sp, dash)
    return sp

def rect(slide, x, y, w, h, fill=None, fa=None, line=None, la=None, lw=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(max(w, 2)), E(max(h, 2)))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = rgb(line); sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    _alpha(sp, fa, la)
    return sp

def oval(slide, x, y, w, h, fill, fa=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, E(x), E(y), E(max(w, 2)), E(max(h, 2)))
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill); sp.line.fill.background(); sp.shadow.inherit = False
    _alpha(sp, fa)
    return sp

def textbox(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_pt=None):
    """paras = list of lines; each line = list of (text, size_pt, color, bold) runs."""
    tb = slide.shapes.add_textbox(E(x), E(y), E(max(w, 8)), E(max(h, 8))); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_pt: p.line_spacing = Pt(line_pt)
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt; set_run(r, size, color, bold)
    return tb

def chip(slide, x, y, w, h, text, size, fill, fg, fa=None, rad=0.35):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, E(x), E(y), E(w), E(h))
    try: sp.adjustments[0] = rad
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill); sp.line.fill.background(); sp.shadow.inherit = False
    if fa is not None: _alpha(sp, fa)
    tf = sp.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; set_run(r, size, fg, True)
    return sp

def picture(slide, path, x, y, w, h):
    import os
    if path and os.path.exists(path):
        return slide.shapes.add_picture(path, E(x), E(y), E(w), E(h))

def bg_fill(slide, color):
    return rect(slide, 0, 0, EMU_W / PXE, EMU_H / PXE, fill=color)

def bg_image(slide, path):
    return slide.shapes.add_picture(path, 0, 0, Emu(EMU_W), Emu(EMU_H))

def save(prs, path):
    prs.save(path)
    xml = zipfile.ZipFile(path).read('ppt/slides/slide1.xml').decode()
    n = len([t for t in re.findall(r'<a:t>([^<]*)</a:t>', xml) if t.strip()])
    print(f'saved {path}')
    print(f'  editable_text_runs={n}  shapes={xml.count("<p:sp>")}  pics={xml.count("<p:pic>")}')
    return dict(text=n, shapes=xml.count('<p:sp>'), pics=xml.count('<p:pic>'))
